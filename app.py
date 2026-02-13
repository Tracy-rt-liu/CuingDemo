import streamlit as st
import time
import random
import pandas as pd
import altair as alt
import io
import re
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

st.set_page_config(layout="wide", page_title="Cuing Agent")

# --- Custom CSS ---
st.markdown("""
<style>
[data-testid="stHorizontalBlock"] { align-items: center; }
.download-row { display: flex; align-items: center; gap: 8px; margin-top: -15px; margin-bottom: 25px; }
.download-label { font-weight: 500; font-size: 0.85rem; color: white; background: #333; padding: 2px 8px; border-radius: 4px; margin-right: 4px; text-transform: uppercase; letter-spacing: 0.5px; }
div.stDownloadButton > button { padding: 4px 10px !important; height: auto !important; min-height: unset !important; border: 1px solid #ddd !important; background: #fff !important; font-size: 1.2rem !important; border-radius: 4px !important; }
div.stDownloadButton > button:hover { border-color: #1f77b4 !important; background: #f0f7ff !important; }
.next-steps-container { border-top: 1px solid #eee; margin-top: 30px; padding-top: 20px; }
.next-steps-title { font-weight: 600; color: #333; margin-bottom: 12px; }
.summary-container { background: #fdfdfd; border: 1px solid #eef0f2; border-radius: 8px; padding: 24px; margin-top: 20px; box-shadow: 0 2px 4px rgba(0,0,0,0.02); width: 100%; }
.summary-section-title { font-size: 1.2rem; font-weight: 700; color: #1f77b4; margin-bottom: 16px; border-bottom: 2px solid #f0f0f0; padding-bottom: 8px; }
.summary-item-header { font-weight: 600; color: #333; margin-top: 12px; margin-bottom: 4px; }
.summary-item-text { color: #555; margin-bottom: 16px; line-height: 1.5; }
.action-cat { font-weight: 600; color: #d62728; margin-top: 14px; margin-bottom: 6px; font-size: 0.9rem; text-transform: uppercase; }
.future-analysis { background: #f0f7ff; border-left: 4px solid #1f77b4; padding: 12px; margin-top: 20px; font-size: 0.9rem; color: #444; }

/* Icon Colors */
.pdf-dl button { color: #d62728 !important; border-color: #ffcccc !important; background-color: #fff5f5 !important; }
.pdf-dl button:hover { background-color: #ffecec !important; border-color: #d62728 !important; }
.xls-dl button { color: #2ca02c !important; border-color: #ccffcc !important; background-color: #f5fff5 !important; }
.xls-dl button:hover { background-color: #e6ff e6 !important; border-color: #2ca02c !important; }
.ppt-dl button { color: #e67e22 !important; border-color: #ffe5cc !important; background-color: #fffaf5 !important; }
.ppt-dl button:hover { background-color: #fff0e0 !important; border-color: #e67e22 !important; }
</style>
""", unsafe_allow_html=True)
st.markdown("", unsafe_allow_html=True) # Spacer

# --- Helper Functions for Exports ---
def strip_emojis(text):
    if not text:
        return ""
    # Remove non-ASCII characters (emojis) for FPDF compatibility
    return re.sub(r'[^\x00-\x7F]+', '', text)

def to_excel(df):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='Sheet1')
    return output.getvalue()

def to_pdf(title, df, insights):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", "B", 16)
    pdf.cell(0, 10, strip_emojis(title), ln=True)
    pdf.ln(5)
    pdf.set_font("helvetica", "", 12)
    pdf.multi_cell(0, 10, strip_emojis(insights))
    pdf.ln(10)
    
    # Table headers
    pdf.set_font("helvetica", "B", 10)
    col_width = pdf.epw / len(df.columns)
    for col in df.columns:
        pdf.cell(col_width, 10, strip_emojis(str(col)), border=1)
    pdf.ln()
    
    # Table rows
    pdf.set_font("helvetica", "", 10)
    for i in range(len(df)):
        for col in df.columns:
            val = df.iloc[i][col]
            pdf.cell(col_width, 10, strip_emojis(str(val)), border=1)
        pdf.ln()
    return bytes(pdf.output())

def to_ppt(title, df, insights):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = strip_emojis(title)
    
    # Insights box
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.text = strip_emojis(insights)
    
    # Add a table for data
    rows, cols = len(df) + 1, len(df.columns)
    left, top, width, height = Inches(0.5), Inches(2.5), Inches(9), Inches(4)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    for i, col in enumerate(df.columns):
        table.cell(0, i).text = strip_emojis(str(col))
    
    # Rows
    for r_idx, row in df.iterrows():
        for c_idx, val in enumerate(row):
            table.cell(r_idx + 1, c_idx).text = strip_emojis(str(val))
            
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

def render_executive_summary():
    st.markdown("<div class='summary-container'>", unsafe_allow_html=True)
    
    st.markdown("<div class='summary-section-title'>📊 Executive Summary</div>", unsafe_allow_html=True)
    
    # Item 1
    st.markdown("<p class='summary-item-header'>1. The OEE decline reflects a sustained performance degradation rather than normal operating variability.</p>", unsafe_allow_html=True)
    st.markdown("<p class='summary-item-text'>The 10.4 percentage point reduction is driven by measurable shifts in availability and quality, indicating a systemic issue rather than statistical fluctuation.</p>", unsafe_allow_html=True)
    
    # Item 2
    st.markdown("<p class='summary-item-header'>2. Line 3 reliability deterioration is the primary operational constraint.</p>", unsafe_allow_html=True)
    st.markdown("<p class='summary-item-text'>Elevated unplanned downtime, reduced MTBF, and near-capacity utilization identify this asset as the bottleneck driving availability loss and throughput instability.</p>", unsafe_allow_html=True)
    
    # Item 3
    st.markdown("<p class='summary-item-header'>3. The scrap increase has become financially material and requires intervention.</p>", unsafe_allow_html=True)
    st.markdown("<p class='summary-item-text'>The 2.7% rise in scrap equates to approximately $243K in annualized margin exposure, with an upward trend over the past six weeks.</p>", unsafe_allow_html=True)
    
    st.markdown("<div class='summary-section-title' style='margin-top: 30px;'>🚀 Recommended Actions</div>", unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown("<p class='action-cat'>Immediate (0–2 weeks)</p>", unsafe_allow_html=True)
        st.markdown("""
        - Conduct focused maintenance audit on Line 3
        - Increase daily scrap tracking by defect category
        - Separate new hires to shadow shifts with highest yield stability
        """)
    with col2:
        st.markdown("<p class='action-cat'>Near-Term (30 days)</p>", unsafe_allow_html=True)
        st.markdown("""
        - Implement daily OEE stand-up dashboard (shift-level)
        - Track individual operator first-pass yield
        - Conduct SMED review on micro-stoppages
        """)
    with col3:
        st.markdown("<p class='action-cat'>Structural</p>", unsafe_allow_html=True)
        st.markdown("""
        - CapEx case for equipment refurbishment
        - Install predictive maintenance sensor on press
        - Formalize ramp KPI: "Time-to-95% Productivity"
        """)
        
    st.markdown("""
    <div class='future-analysis'>
    <strong>I can add additional analysis below if helpful:</strong><br>
    • Monte Carlo simulation for projected yield under different scenarios<br>
    • Sensitivity analysis: “If availability improves 3pp → output impact = X”<br>
    • Automated financial impact estimator<br>
    • Benchmark comparison vs industry quartiles
    </div>
    """, unsafe_allow_html=True)
    
    st.markdown("</div>", unsafe_allow_html=True)

st.title("🏭 Cuing Agent")

# --- Session State ---
if "messages" not in st.session_state:
    st.session_state.messages = []
if "streaming_done" not in st.session_state:
    st.session_state.streaming_done = True
if "answers" not in st.session_state:
    st.session_state.answers = {}
if "step" not in st.session_state:
    st.session_state.step = 0

# --- Response Definitions ---
STEP_0_BLOCKS = [
    {
        "type": "text",
        "content": (
            "To diagnose the yield decline, I'd like to validate four potential drivers:\n\n"
            "| Metric | Last Week | 4-Week Avg | Delta |\n"
            "| --- | --- | --- | --- |\n"
            "| First Pass Yield | 88.2% | 92.5% | -4.3pp |\n"
            "| Throughput / Worker | 13 units/day | 15 units/day | -13% |\n"
            "| OEE | 71% | 78% | -7pp |\n"
            "| Scrap Rate | 6.5% | 3.8% | +2.7pp |\n"
            "| Rework Hours | +22% | Baseline | ↑ |\n\n"
            "**Labor ramp & productivity**  \n"
            "**Equipment performance (OEE)**  \n"
            "**Material quality**  \n"
            "**Process bottlenecks / cycle time shifts**\n\n"
            "Based on your SAP + MES integration, here's what I see for "
            "last week vs trailing 4-week average:"
        ),
    },
    {
        "type": "text",
        "content": (
            "\n---\n\n"
            "**🔧 Labor Changes Identified**\n"
            "- 3 new workers onboarded 4 weeks ago\n"
            "- Standard training ramp: 4 weeks\n"
            "- Current effective productivity of new workers estimated "
            "at 70–80% of experienced operator\n\n"
            "**Please confirm:**"
        ),
    },
    {
        "type": "questions",
        "items": [
            {"id": "ramp_assumption", "text": "Is the 4-week ramp assumption accurate?"},
            {"id": "new_hires_line", "text": "Are the new hires assigned to the affected line?"},
            {"id": "scrap_shifts", "text": "Has scrap increased on specific shifts?"},
        ],
    },
    {
        "type": "text",
        "content": (
            "\n---\n\n"
            "**⚙️ Equipment Observations**\n"
            "- Line 3 stamping press last major overhaul: March 2021\n"
            "- Recommended refurbishment cycle: 18–24 months\n"
            "- Unplanned downtime last week: +18%\n"
            "- Minor stoppages: +27%\n\n"
            "**Can you confirm:**"
        ),
    },
    {
        "type": "questions",
        "items": [
            {"id": "micro_stoppages", "text": "Have you experienced increased micro-stoppages?"},
            {"id": "maintenance_deferrals", "text": "Any recent maintenance deferrals?"},
        ],
    },
]

STEP_1_BLOCKS = [
    {
        "type": "text",
        "content": (
            "Thank you for confirming the inputs. Below is a structured "
            "assessment of impact by driver, ranked by materiality."
        ),
    },
    {
        "type": "text",
        "content": (
            "\n---\n\n"
            "**A. Labor Ramp Impact – Low to Moderate**\n\n"
            "**System Data:**\n"
            "- 20 operators total\n"
            "- 3 new operators\n"
            "- Experienced productivity: 15 units/day\n"
            "- New operator productivity during ramp: 75%\n\n"
            "**Effective Output Calculation**\n\n"
            "Experienced output:\n"
            "17 workers × 15 units = **255 units/day**\n\n"
            "New workers:\n"
            "3 × (15 × 75%) = **33.75 units/day**\n\n"
            "Total expected output:\n"
            "255 + 33.75 = **288.75 units/day**\n\n"
            "If fully ramped:\n"
            "20 × 15 = **300 units/day**\n\n"
            "Ramp drag = **11.25 units/day** (3.75% total output loss)\n\n"
            "This explains some throughput loss but **not** a 4.3pp yield drop.\n\n"
            "**Conclusion:** Labor ramp is contributory, not primary driver."
        ),
    },
    {
        "type": "text",
        "content": (
            "\n---\n\n"
            "**B. OEE Decomposition – Primary Driver**\n\n"
            "OEE = Availability × Performance × Quality\n\n"
            "**From system data:**\n"
            "- Availability dropped from 90% → 84%\n"
            "- Performance dropped from 95% → 92%\n"
            "- Quality dropped from 98% → 95%\n\n"
            "New OEE: 0.84 × 0.92 × 0.95 = **73.4%**\n"
            "Previous OEE: 0.90 × 0.95 × 0.98 = **83.8%**\n\n"
            "That's a **10.4pp OEE deterioration**. The decline is primarily driven "
            "by reduced Availability and Quality. The magnitude of OEE loss "
            "significantly exceeds the labor ramp effect and aligns with observed "
            "yield and scrap increases."
        ),
    },
    {
        "type": "text",
        "content": (
            "\n---\n\n"
            "**C. Scrap Financial Impact – Economic Impact**\n\n"
            "**System Data:**\n"
            "- 300 units/day baseline\n"
            "- Contribution margin per unit: \\$120\n\n"
            "**Analysis:**\n"
            "- Scrap increase: (6.5% – 3.8%) × 300 = **8.1 additional scrap units/day**\n"
            "- Financial impact: 8.1 × \\$120 = **\\$972/day**\n"
            "- Annualized (250 days): **\\$243,000 impact**\n\n"
            "Scrap increase is economically material and consistent with the observed "
            "quality degradation within OEE."
        ),
    },
    {
        "type": "text",
        "content": (
            "\n---\n\n"
            "Please confirm whether you would like me to incorporate more variables "
            "into the next diagnostic view, or proceed with visualization focused on "
            "equipment-driven yield deterioration."
        ),
    },
]

STEP_2_BLOCKS = [
    {
        "type": "chart",
        "title": "1️⃣ OEE Waterfall (Week-over-Week)",
        "chart_type": "oee_waterfall",
        "commentary": (
            "**58%** of deterioration driven by availability loss\n\n"
            "**29%** driven by performance slowdown\n\n"
            "**13%** driven by quality decline\n\n"
            "**Confirms equipment instability as dominant driver**"
        ),
    },
    {
        "type": "chart",
        "title": "2️⃣ Constrained Asset Performance",
        "chart_type": "asset_performance",
        "commentary": (
            "Asset running near theoretical capacity\n\n"
            "Reliability has deteriorated materially\n\n"
            "Breakdown frequency increased ~4×\n\n"
            "**Clear constraint machine impacting system flow**"
        ),
    },
    {
        "type": "chart",
        "title": "3️⃣ Scrap & Margin Impact Trend",
        "chart_type": "scrap_trend",
        "commentary": (
            "**Approximate cumulative impact to date:** ≈ $32–38K realized in last 6 weeks"
        ),
    },
]

FALLBACK_BLOCKS = [
    {
        "type": "text",
        "content": "Thank you for the additional context. Let me analyze this further "
        "and get back to you with updated recommendations.",
    }
]


def get_response_blocks(step):
    if step == 0:
        return STEP_0_BLOCKS
    if step == 1:
        return STEP_1_BLOCKS
    if step == 2:
        return STEP_2_BLOCKS
    return FALLBACK_BLOCKS


# --- Thinking animation ---
THINKING_PHASES = {
    0: [
        ("Querying SAP + MES data...", 1.2),
        ("Analyzing yield trends...", 1.0),
        ("Cross-referencing equipment logs...", 0.8),
        ("Preparing diagnostic summary...", 0.6),
    ],
    1: [
        ("Retrieving operator headcount and new hire records...", 0.4),
        ("Calculating effective productivity for new vs experienced workers...", 0.35),
        ("Aggregating total expected output and ramp drag impact...", 0.3),
        ("Comparing projected labor output vs actual yield...", 0.3),
        ("Flagging labor contribution to yield deviation...", 0.25),
        ("Pulling line-level OEE metrics from MES logs...", 0.35),
        ("Decomposing Availability, Performance, and Quality components...", 0.3),
        ("Computing week-over-week OEE change...", 0.25),
        ("Evaluating relative impact of each OEE component on total output...", 0.3),
        ("Extracting daily scrap counts and defect categories...", 0.3),
        ("Calculating incremental units lost to scrap...", 0.25),
        ("Converting scrap units into financial impact based on contribution margin...", 0.3),
        ("Annualizing margin exposure from weekly scrap data...", 0.25),
        ("Prioritizing drivers by combined operational and economic impact...", 0.3),
        ("Preparing structured assessment summary for user confirmation...", 0.35),
    ],
}
THINKING_PHASES[2] = [
    ("Creating data visualization...", 5.0),
]


def show_thinking(step=0):
    """Display an animated thinking state before the response."""
    phases = THINKING_PHASES.get(step, THINKING_PHASES[0])
    with st.status("🔍 Analyzing system data...", expanded=True) as status:
        for label, duration in phases:
            st.write(label)
            time.sleep(duration)
        status.update(label="✅ Analysis complete", state="complete", expanded=False)
    time.sleep(0.3)


# --- Streaming helpers ---
def text_generator(text, speed=0.03):
    """Yield text word-by-word, batching table rows as one chunk."""
    lines = text.split("\n")
    i = 0
    while i < len(lines):
        if lines[i].strip().startswith("|"):
            # Render the entire table at once
            table = ""
            while i < len(lines) and lines[i].strip().startswith("|"):
                table += lines[i] + "\n"
                i += 1
            yield table
            time.sleep(0.4)
        else:
            line = lines[i]
            i += 1
            if line.strip():
                # Stream word-by-word for natural feel
                words = line.split(" ")
                for j, word in enumerate(words):
                    yield word + (" " if j < len(words) - 1 else "")
                    time.sleep(speed + random.uniform(0, 0.02))
                yield "\n"
            else:
                yield "\n"


def render_questions(items):
    """Render interactive Yes / No pill-buttons for each question."""
    for item in items:
        qid = item["id"]
        answer = st.session_state.answers.get(qid)

        cols = st.columns([5, 0.8, 0.8, 3.4])
        with cols[0]:
            st.markdown(f"&ensp;{item['text']}")
        with cols[1]:
            yes_type = "primary" if answer == "yes" else "secondary"
            yes_label = "✓ Yes" if answer == "yes" else "Yes"
            if st.button(yes_label, key=f"{qid}_yes", type=yes_type, use_container_width=True):
                st.session_state.answers[qid] = "yes"
                st.rerun()
        with cols[2]:
            no_type = "primary" if answer == "no" else "secondary"
            no_label = "✗ No" if answer == "no" else "No"
            if st.button(no_label, key=f"{qid}_no", type=no_type, use_container_width=True):
                st.session_state.answers[qid] = "no"
                st.rerun()


def render_chart_oee_waterfall():
    data = pd.DataFrame([
        {"label": "Previous OEE", "amount": 83.8, "type": "total"},
        {"label": "Availability", "amount": -6, "type": "delta"},
        {"label": "Performance", "amount": -3, "type": "delta"},
        {"label": "Quality", "amount": -1.4, "type": "delta"},
        {"label": "New OEE", "amount": 73.4, "type": "total"},
    ])
    
    # Pre-calculate start and end values for waterfall bars
    data["end"] = data["amount"].cumsum()
    data["start"] = data["end"].shift(1).fillna(0)
    
    # Adjust start/end for 'total' bars (they start at 0)
    data.loc[data["type"] == "total", "start"] = 0
    data.loc[data["type"] == "total", "end"] = data["amount"]
    
    # Logic for 'delta' bars: if amount is negative, start is higher than end
    # We need to ensure 'start' is the previous cumulative sum and 'end' is new cumulative sum.
    # Actually, simpler: just track the running total.
    # Previous: 0 -> 83.8
    # Avail: 83.8 -> 77.8
    # Perf: 77.8 -> 74.8
    # Qual: 74.8 -> 73.4
    # New: 0 -> 73.4
    
    # Correct calculation:
    running_total = 0
    starts = []
    ends = []
    for _, row in data.iterrows():
        if row["type"] == "total":
            starts.append(0)
            ends.append(row["amount"])
            running_total = row["amount"]
        else:
            starts.append(running_total)
            running_total += row["amount"]
            ends.append(running_total)
    
    data["start"] = starts
    data["end"] = ends
    data["color"] = data["type"].apply(lambda x: "#1f77b4" if x == "total" else "#ff7f0e")
    
    c = (
        alt.Chart(data)
        .mark_bar()
        .encode(
            x=alt.X("label", sort=None, title=None, axis=alt.Axis(labelAngle=0)),
            y=alt.Y("start", title="OEE %"),
            y2="end",
            color=alt.Color("color", scale=None),
            tooltip=["label", "amount", "end"],
        )
        .properties(height=300)
    )
    st.altair_chart(c, use_container_width=True)
    return data


def render_chart_asset_performance():
    data = pd.DataFrame([
        {"metric": "Capacity Utilization", "value": 94, "unit": "%"},
        {"metric": "Downtime", "value": 18, "unit": "%"},
        {"metric": "MTBF", "value": 22, "unit": "hours"},
    ])
    
    c = (
        alt.Chart(data)
        .mark_bar()
        .encode(
            x=alt.X("metric", sort=None, title=None, axis=alt.Axis(labelAngle=0)),
            y=alt.Y("value", title="Value"),
            color=alt.value("#1f77b4"),
            tooltip=["metric", "value", "unit"],
        )
        .properties(height=300)
    )
    st.altair_chart(c, use_container_width=True)
    return data


def render_chart_scrap_trend(key_prefix=""):
    # Layout for numeric input
    c1, c2 = st.columns([2, 1])
    with c1:
        st.markdown("<div style='padding-top: 10px;'>Select duration (Weeks):</div>", unsafe_allow_html=True)
    with c2:
        num_weeks = st.number_input(
            "Select duration (Weeks):",
            min_value=1,
            max_value=12,
            value=6,
            step=1,
            key=f"scrap_view_input_{key_prefix}",
            label_visibility="collapsed"
        )
    
    weeks = num_weeks
    
    # Data
    data = pd.DataFrame({
        "Week": range(1, 9),
        "Scrap %": [3.8, 4.1, 4.4, 4.8, 5.7, 6.5, 6.2, 6.0],
        "Margin Loss ($k)": [0, 3.6, 7.2, 14.4, 21.6, 29.7, 26.4, 24.0],
    })
    data = data[data["Week"] <= weeks]
    
    base = alt.Chart(data).encode(x=alt.X("Week:O", axis=alt.Axis(labelAngle=0)))

    line_scrap = base.mark_line(color="#ff7f0e").encode(
        y=alt.Y("Scrap %:Q", title="Scrap %", axis=alt.Axis(titleColor="#ff7f0e", orient='left')),
        tooltip=["Week", "Scrap %"]
    )
    
    line_margin = base.mark_line(color="#1f77b4", strokeDash=[5, 5]).encode(
        y=alt.Y("Margin Loss ($k):Q", title="Weekly Margin Loss ($k)", axis=alt.Axis(titleColor="#1f77b4", orient='right')),
        tooltip=["Week", "Margin Loss ($k)"]
    )
    
    # Baseline annotation (3.8%) - share left axis scale but suppress redundant label
    rule = alt.Chart(pd.DataFrame({"Scrap %": [3.8]})).mark_rule(strokeDash=[2, 2], color="gray").encode(
        y=alt.Y("Scrap %:Q", axis=None)
    )

    c = alt.layer(line_scrap, line_margin, rule).resolve_scale(y="independent").properties(height=350)
    
    st.altair_chart(c, use_container_width=True)
    return data


def render_blocks(blocks, streaming=False, step=0, key_prefix=""):
    """Render response blocks; stream text progressively when streaming=True."""
    if streaming:
        show_thinking(step)
    for block in blocks:
        if block["type"] == "text":
            if streaming:
                st.write_stream(text_generator(block["content"]))
            else:
                st.markdown(block["content"])
        elif block["type"] == "questions":
            if streaming:
                time.sleep(0.3)
            render_questions(block["items"])
        elif block["type"] == "chart":
            st.subheader(block["title"])
            df = None
            if block["chart_type"] == "oee_waterfall":
                df = render_chart_oee_waterfall()
            elif block["chart_type"] == "asset_performance":
                df = render_chart_asset_performance()
            elif block["chart_type"] == "scrap_trend":
                df = render_chart_scrap_trend(key_prefix)
            
            # Download Section
            if df is not None:
                cols = st.columns([1.5, 0.6, 0.6, 0.6, 6.7])
                with cols[0]:
                    st.markdown("<div style='padding-top: 6px;'><span class='download-label'>Download:</span></div>", unsafe_allow_html=True)
                with cols[1]:
                    st.markdown('<div class="pdf-dl">', unsafe_allow_html=True)
                    st.download_button("📄", data=to_pdf(block["title"], df, block["commentary"]), file_name=f"{block['chart_type']}.pdf", key=f"dl_pdf_{block['chart_type']}_{key_prefix}", help="Download PDF Report")
                    st.markdown('</div>', unsafe_allow_html=True)
                with cols[2]:
                    st.markdown('<div class="xls-dl">', unsafe_allow_html=True)
                    st.download_button("📊", data=to_excel(df), file_name=f"{block['chart_type']}.xlsx", key=f"dl_xls_{block['chart_type']}_{key_prefix}", help="Download Excel Data")
                    st.markdown('</div>', unsafe_allow_html=True)
                with cols[3]:
                    st.markdown('<div class="ppt-dl">', unsafe_allow_html=True)
                    st.download_button("📽️", data=to_ppt(block["title"], df, block["commentary"]), file_name=f"{block['chart_type']}.pptx", key=f"dl_ppt_{block['chart_type']}_{key_prefix}", help="Download PPT Slides")
                    st.markdown('</div>', unsafe_allow_html=True)
                st.markdown("<div style='margin-bottom: 10px;'></div>", unsafe_allow_html=True)

            if streaming:
                st.write_stream(text_generator(block["commentary"]))
            else:
                st.markdown(block["commentary"])
            st.divider()


# --- Display Messages ---
for i, msg in enumerate(st.session_state.messages):
    is_last = i == len(st.session_state.messages) - 1

    with st.chat_message(msg["role"]):
        if msg["role"] == "user":
            st.markdown(msg["content"])
        else:
            blocks = get_response_blocks(msg.get("step", 0))
            should_stream = is_last and not st.session_state.streaming_done
            if should_stream:
                st.session_state.streaming_done = True  # prevent re-stream on rerun
            render_blocks(blocks, streaming=should_stream, step=msg.get("step", 0), key_prefix=str(i))
            
            # Next Steps Section (Only after Step 2: Analysis)
            if msg.get("step", 0) == 2:
                st.markdown("<div class='next-steps-container'></div>", unsafe_allow_html=True)
                st.markdown("<p class='next-steps-title'>Next, do you want me to:</p>", unsafe_allow_html=True)
                
                c1, c2 = st.columns(2)
                with c1:
                    if st.checkbox("1. Make edits to these charts", key=f"next_edits_{i}"):
                        st.info("I'm ready to help you refine these visualizations. What changes would you like to see?")
                with c2:
                    show_summary = st.checkbox("2. Output Executive Summary and Recommended Actions", key=f"next_summary_{i}")
                
                if show_summary:
                    render_executive_summary()

# --- Chat Input ---
user_input = st.chat_input("Ask about yield performance...")

if user_input:
    current_step = st.session_state.step
    st.session_state.messages.append({"role": "user", "content": user_input})
    st.session_state.messages.append({"role": "assistant", "step": current_step})
    st.session_state.step += 1
    st.session_state.streaming_done = False
    st.rerun()
